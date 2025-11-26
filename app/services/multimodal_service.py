"""
멀티모달 서비스
이미지 추출 및 텍스트 변환
"""

import logging
from typing import List, Dict, Optional
from uuid import UUID
import base64
import io
from PIL import Image
import fitz  # PyMuPDF
import aiohttp

from app.model.Document import Document, FileType

logger = logging.getLogger(__name__)


class MultimodalService:
    """멀티모달(텍스트 + 이미지) 처리 서비스"""
    
    def __init__(self):
        pass
    
    async def extract_content_with_images(
        self, 
        document: Document,
        content: str
    ) -> Dict[str, any]:
        """
        문서에서 텍스트와 이미지를 추출
        
        Args:
            document: Document 엔티티
            content: 문서 텍스트 내용
            
        Returns:
            {
                "text": str,
                "images": List[Dict] - [{data: base64, page: int, caption: str}]
            }
        """
        images = []
        text = content
        
        # PDF 파일인 경우 텍스트 + 이미지 추출
        if document.content_type and 'pdf' in document.content_type.lower():
            if document.type == FileType.FILE:
                result = await self._extract_from_pdf_s3(document.value)
                text = result.get("text", content)
                images = result.get("images", [])
            elif document.type == FileType.URL:
                result = await self._extract_from_pdf_url(document.value)
                text = result.get("text", content)
                images = result.get("images", [])
        
        # PPT 파일인 경우 텍스트 + 이미지 추출
        elif document.content_type and ('presentation' in document.content_type.lower() or 
                                        'powerpoint' in document.content_type.lower() or
                                        document.value.endswith(('.ppt', '.pptx'))):
            if document.type == FileType.FILE:
                result = await self._extract_from_ppt_s3(document.value)
                text = result.get("text", content)
                images = result.get("images", [])
            elif document.type == FileType.URL:
                result = await self._extract_from_ppt_url(document.value)
                text = result.get("text", content)
                images = result.get("images", [])
        
        # 이미지 파일인 경우
        elif document.content_type and document.content_type.startswith('image/'):
            if document.type == FileType.FILE:
                image_data = await self._get_image_from_s3(document.value)
                if image_data:
                    images = [{
                        "data": image_data,
                        "page": 1,
                        "caption": document.title or "Uploaded image"
                    }]
            elif document.type == FileType.URL:
                image_data = await self._get_image_from_url(document.value)
                if image_data:
                    images = [{
                        "data": image_data,
                        "page": 1,
                        "caption": document.title or "Image from URL"
                    }]
        
        return {
            "text": text,
            "images": images
        }
    
    async def _extract_from_pdf_s3(self, s3_key: str) -> Dict:
        """S3의 PDF에서 텍스트와 이미지 추출"""
        try:
            import boto3
            from app.config.config import settings
            
            # S3에서 PDF 다운로드
            s3_client = boto3.client(
                's3',
                region_name=settings.aws_region,
                aws_access_key_id=settings.aws_access_key_id,
                aws_secret_access_key=settings.aws_secret_access_key
            ) if settings.aws_access_key_id else boto3.client('s3', region_name=settings.aws_region)
            
            response = s3_client.get_object(Bucket=settings.aws_s3_bucket, Key=s3_key)
            pdf_bytes = response['Body'].read()
            
            return self._extract_from_pdf_bytes(pdf_bytes)
            
        except Exception as e:
            logger.error(f"Failed to extract from PDF S3: {e}")
            return {"text": "", "images": []}
    
    async def _extract_from_pdf_url(self, url: str) -> Dict:
        """URL의 PDF에서 텍스트와 이미지 추출"""
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=60) as response:
                    if response.status == 200:
                        pdf_bytes = await response.read()
                        return self._extract_from_pdf_bytes(pdf_bytes)
            return {"text": "", "images": []}
        except Exception as e:
            logger.error(f"Failed to extract from PDF URL: {e}")
            return {"text": "", "images": []}
    
    def _extract_from_pdf_bytes(self, pdf_bytes: bytes) -> Dict:
        """PDF 바이트에서 텍스트와 이미지 추출"""
        images = []
        text_parts = []
        
        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 텍스트 추출
                page_text = page.get_text()
                if page_text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                
                # 이미지 추출
                image_list = page.get_images(full=True)
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # 이미지를 base64로 인코딩
                    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    
                    # 이미지가 너무 작으면 스킵 (아이콘 등)
                    try:
                        img_obj = Image.open(io.BytesIO(image_bytes))
                        if img_obj.width < 100 or img_obj.height < 100:
                            continue
                    except:
                        continue
                    
                    images.append({
                        "data": image_base64,
                        "page": page_num + 1,
                        "caption": f"Page {page_num + 1}, Image {img_index + 1}",
                        "format": base_image["ext"]
                    })
            
            doc.close()
            text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(text_parts)} pages text and {len(images)} images from PDF")
            
        except Exception as e:
            logger.error(f"Failed to extract from PDF bytes: {e}")
            text = ""
        
        return {
            "text": text,
            "images": images
        }
    
    async def _extract_from_ppt_s3(self, s3_key: str) -> Dict:
        """S3의 PPT에서 텍스트와 이미지 추출"""
        try:
            import boto3
            from app.config.config import settings
            
            # S3에서 PPT 다운로드
            s3_client = boto3.client(
                's3',
                region_name=settings.aws_region,
                aws_access_key_id=settings.aws_access_key_id,
                aws_secret_access_key=settings.aws_secret_access_key
            ) if settings.aws_access_key_id else boto3.client('s3', region_name=settings.aws_region)
            
            response = s3_client.get_object(Bucket=settings.aws_s3_bucket, Key=s3_key)
            ppt_bytes = response['Body'].read()
         
            return self._extract_from_ppt_bytes(ppt_bytes)
            
        except Exception as e:
            logger.error(f"Failed to extract from PPT S3: {e}")
            return {"text": "", "images": []}
    
    async def _extract_from_ppt_url(self, url: str) -> Dict:
        """URL의 PPT에서 텍스트와 이미지 추출"""
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=60) as response:
                    if response.status == 200:
                        ppt_bytes = await response.read()
                        return self._extract_from_ppt_bytes(ppt_bytes)
            return {"text": "", "images": []}
        except Exception as e:
            logger.error(f"Failed to extract from PPT URL: {e}")
            return {"text": "", "images": []}
    
    def _extract_from_ppt_bytes(self, ppt_bytes: bytes) -> Dict:
        """PPT 바이트에서 텍스트와 이미지 추출"""
        from pptx import Presentation
        
        images = []
        text_parts = []
        
        try:
            # BytesIO로 감싸서 Presentation 객체 생성
            ppt_stream = io.BytesIO(ppt_bytes)
            prs = Presentation(ppt_stream)
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                
                # 슬라이드의 모든 shape에서 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # 이미지 추출
                    if shape.shape_type == 13:  # Picture type
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            
                            # 이미지 크기 확인
                            img_obj = Image.open(io.BytesIO(image_bytes))
                            if img_obj.width < 100 or img_obj.height < 100:
                                continue
                            
                            # base64 인코딩
                            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                            
                            images.append({
                                "data": image_base64,
                                "page": slide_num,
                                "caption": f"Slide {slide_num}, Image",
                                "format": image.ext
                            })
                        except Exception as e:
                            logger.debug(f"Failed to extract image from slide {slide_num}: {e}")
                            continue
                
                # 슬라이드 텍스트 추가
                if slide_text:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
            
            text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(text_parts)} slides text and {len(images)} images from PPT")
            
        except Exception as e:
            logger.error(f"Failed to extract from PPT bytes: {e}")
            text = ""
        
        return {
            "text": text,
            "images": images
        }
    
    async def _get_image_from_s3(self, s3_key: str) -> Optional[str]:
        """S3에서 이미지 가져오기 (base64)"""
        try:
            import boto3
            from app.config.config import settings
            
            s3_client = boto3.client(
                's3',
                region_name=settings.aws_region,
                aws_access_key_id=settings.aws_access_key_id,
                aws_secret_access_key=settings.aws_secret_access_key
            ) if settings.aws_access_key_id else boto3.client('s3', region_name=settings.aws_region)
            
            response = s3_client.get_object(Bucket=settings.aws_s3_bucket, Key=s3_key)
            image_bytes = response['Body'].read()
            
            return base64.b64encode(image_bytes).decode('utf-8')
            
        except Exception as e:
            logger.error(f"Failed to get image from S3: {e}")
            return None
    
    async def _get_image_from_url(self, url: str) -> Optional[str]:
        """URL에서 이미지 가져오기 (base64)"""
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=30) as response:
                    if response.status == 200:
                        image_bytes = await response.read()
                        return base64.b64encode(image_bytes).decode('utf-8')
            return None
        except Exception as e:
            logger.error(f"Failed to get image from URL: {e}")
            return None


# 싱글톤 인스턴스
_multimodal_service: Optional[MultimodalService] = None


def get_multimodal_service() -> MultimodalService:
    """MultimodalService 싱글톤 인스턴스 반환"""
    global _multimodal_service
    if _multimodal_service is None:
        _multimodal_service = MultimodalService()
    return _multimodal_service
